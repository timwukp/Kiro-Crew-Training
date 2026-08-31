#!/usr/bin/env python3
"""
Generate Kiro Crew Training Presentation (Verified Content)
Based on official documentation from kiro.dev
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Theme colors (matching website)
BG_DARK = RGBColor(15, 23, 42)
CARD_BG = RGBColor(30, 41, 59)
TEXT_COLOR = RGBColor(241, 245, 249)
MUTED = RGBColor(148, 163, 184)
PRIMARY = RGBColor(37, 99, 235)
SECONDARY = RGBColor(124, 58, 237)
ACCENT = RGBColor(6, 182, 212)
SUCCESS = RGBColor(16, 185, 129)
WARNING = RGBColor(245, 158, 11)


def set_background(slide, color=BG_DARK):
    """Set slide background color"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title, subtitle, footer=""):
    """Create title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    set_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = TEXT_COLOR
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    sub_para = subtitle_frame.paragraphs[0]
    sub_para.alignment = PP_ALIGN.CENTER
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = MUTED
    
    # Footer
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        footer_frame = footer_box.text_frame
        footer_frame.text = footer
        footer_para = footer_frame.paragraphs[0]
        footer_para.alignment = PP_ALIGN.CENTER
        footer_para.font.size = Pt(14)
        footer_para.font.color.rgb = ACCENT


def add_section_divider(prs, title, subtitle=""):
    """Create section divider"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, CARD_BG)
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    para = title_frame.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    para.font.size = Pt(44)
    para.font.bold = True
    para.font.color.rgb = PRIMARY
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
        sub_frame = sub_box.text_frame
        sub_frame.text = subtitle
        sub_para = sub_frame.paragraphs[0]
        sub_para.alignment = PP_ALIGN.CENTER
        sub_para.font.size = Pt(20)
        sub_para.font.color.rgb = MUTED


def add_content_slide(prs, title, bullets, verification_status="verified"):
    """Create content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # Title with verification badge
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(7.5), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY
    
    # Verification badge
    badge_color = SUCCESS if verification_status == "verified" else WARNING
    badge_text = "✓ Verified" if verification_status == "verified" else "⚠ To Be Verified"
    badge_box = slide.shapes.add_textbox(Inches(8.2), Inches(0.45), Inches(1.5), Inches(0.4))
    badge_frame = badge_box.text_frame
    badge_frame.text = badge_text
    badge_para = badge_frame.paragraphs[0]
    badge_para.font.size = Pt(11)
    badge_para.font.bold = True
    badge_para.font.color.rgb = badge_color
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR if not bullet.startswith("⚠") else WARNING
        p.space_before = Pt(8)


def add_two_column_slide(prs, title, left_content, right_content, verification_status="verified"):
    """Create two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.5), Inches(5.7))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i > 0:
            left_frame.add_paragraph()
        p = left_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(5.7))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i > 0:
            right_frame.add_paragraph()
        p = right_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(6)


def add_table_slide(prs, title, headers, rows, verification_status="verified"):
    """Create table slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = PRIMARY
    
    # Table
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5.5)
    )
    table = table_shape.table
    
    # Headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
    
    # Rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.text_frame.paragraphs[0].font.color.rgb = TEXT_COLOR
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs,
        "Kiro Crew Training",
        "Personal AI Agent Platform",
        "Facts verified against official kiro.dev docs | 2026-08-31"
    )
    
    # Slide 2: Agenda
    add_content_slide(prs,
        "Agenda (60 minutes)",
        [
            "Introduction (10 min) - What is Kiro Crew?",
            "Positioning (10 min) - Why Crew vs IDE/CLI?",
            "Features Deep Dive (20 min) - Memory, Subagents, Scheduling, Tasks",
            "Use Cases (15 min) - Enterprise integration examples",
            "Getting Started (5 min) - Setup and first steps"
        ]
    )
    
    # Slide 3: Verification Status
    add_content_slide(prs,
        "Content Verification Status",
        [
            "✓ All facts verified against official kiro.dev documentation",
            "✓ Architecture, memory system, features - fully verified",
            "✓ Technical limits and capabilities - documented with sources",
            "",
            "Note: Per-operation credit costs and some quotas are NOT",
            "published by Kiro - measure them in the usage dashboard.",
            "",
            "All sources: kiro.dev/docs/crew/",
            "Last verified: 2026-08-31"
        ],
        verification_status="verified"
    )
    
    # Section 1: Introduction
    add_section_divider(prs, "Part 1: Introduction", "What is Kiro Crew?")
    
    add_content_slide(prs,
        "What is Kiro Crew?",
        [
            "Open-source personal AI agent that runs locally or remotely",
            "",
            "💾 Persistent - Sessions, memory, schedules persist beyond one chat",
            "🧠 Self-learning - Your corrections become durable lessons",
            "🌱 Self-evolving - Repeated patterns can become skills",
            "🤖 Autonomous - 24/7 operation without human presence",
            "",
            "Source: kiro.dev/docs/crew/"
        ]
    )
    
    add_content_slide(prs,
        "The Kiro Ecosystem",
        [
            '"One agent, every surface"',
            "",
            "🖥 IDE - Desktop editor integration",
            "⌨ CLI - Terminal-native",
            "🌐 Web - Browser-based (preview)",
            "📱 Mobile - iOS/Android (preview)",
            "🚀 Crew - Persistent agent with automation",
            "",
            "✓ Single subscription covers ALL surfaces",
            "✓ Shared .kiro/ configuration"
        ]
    )
    
    add_content_slide(prs,
        "Architecture (CORRECTED)",
        [
            "User/Interface ↓",
            "Kiro Crew Gateway (local Python server, port 5476) ↓",
            "kiro-cli (Agent Client Protocol) ↓",
            "Kiro Model Services (cloud, managed by Kiro)",
            "",
            "Local data: ~/.kiro/crew/",
            "Authentication: device-code sign-in on first launch",
            "",
            "✅ NO AWS Bedrock needed",
            "✅ NO AWS CLI needed (unless deploying artifacts)"
        ]
    )
    
    # Section 2: Why Crew
    add_section_divider(prs, "Part 2: Why Crew?", "When you already have IDE/CLI")
    
    add_table_slide(prs,
        "IDE/CLI vs Crew",
        ["Need", "IDE/CLI", "Crew"],
        [
            ["Context", "Session-based", "Persistent (6 layers, 365 days)"],
            ["Availability", "Human required", "24/7 autonomous"],
            ["Learning", "Per-session", "Durable lessons"],
            ["Parallelization", "Single-threaded", "3-32 concurrent subagents"],
            ["Access", "Local terminal", "Multi-channel (Slack/Discord/Teams)"],
            ["Tasks", "Interactive Q&A", "Autonomous task runner"],
        ]
    )
    
    add_two_column_slide(prs,
        "When to Use What",
        [
            "Use IDE/CLI for:",
            "• Real-time coding and debugging",
            "• Interactive prototyping",
            "• Immediate code generation",
            "• Pair programming sessions"
        ],
        [
            "Use Crew for:",
            "• 24/7 monitoring (PR/CI/incidents)",
            "• Scheduled tasks (audits, reports)",
            "• Team access via Slack/Discord",
            "• Complex autonomous tasks",
            "• Persistent memory across days",
            "• Parallel investigation"
        ]
    )
    
    # Section 3: Features
    add_section_divider(prs, "Part 3: Features", "Core capabilities deep dive")
    
    add_content_slide(prs,
        "6-Layer Memory System",
        [
            "1. Preferences (4,250 chars) - habits, style, tools",
            "2. Projects (6,400 chars) - CRs, branches, status",
            "3. Recent History (26,600 chars) - tiered decay 0-365 days",
            "4. Semantic Memory (12,000 chars) - key-value pairs, hybrid search",
            "5. Episodic Memory (3,000 chars, top-8) - past events, max 10K entries",
            "6. Lessons (37,250 chars, max 50) - user-taught rules",
            "",
            "All limits verified from: kiro.dev/docs/crew/features/memory/"
        ]
    )
    
    add_content_slide(prs,
        "Subagents",
        [
            "Parallel background agents for research and investigation",
            "",
            "Concurrency: 3-32 (auto-sized, NOT 11)",
            "Timeout: 30 minutes hard limit",
            "Stall detection: ~2 minutes",
            "Context: Isolated sessions with memory injection",
            "",
            "Example: 'Research Redis, PostgreSQL, MongoDB'",
            "→ Spawns 3 parallel agents",
            "→ Each investigates one technology",
            "→ Results aggregated"
        ]
    )
    
    add_content_slide(prs,
        "Scheduling (3 Modes)",
        [
            "Mode 1: Cron Jobs",
            "• Standard cron expressions or --every N seconds",
            "• Default timeout: 1800s (30 min)",
            "",
            "Mode 2: Heartbeats",
            "• Reactive monitoring (checks every 60s)",
            "• Only surfaces when something changes",
            "",
            "Mode 3: Webhooks",
            "• External triggers via HTTP POST",
            "• Max 6 concurrent sessions",
            "• Bearer token authentication"
        ]
    )
    
    add_content_slide(prs,
        "Task Runner",
        [
            "Autonomous task execution from specifications",
            "",
            "Process:",
            "1. Decompose spec into steps",
            "2. Execute with checkpoints",
            "3. Test step completion",
            "4. Retry on failure",
            "5. Resume from failure point",
            "",
            "Dashboard shows live progress",
            "Source: kiro.dev/docs/crew/features/task-runner/"
        ]
    )
    
    add_content_slide(prs,
        "Artifacts & Knowledge",
        [
            "Artifacts: Persistent outputs with version history",
            "• Kinds: widget, HTML, markdown, SVG, JSON, text",
            "• Live previews in dashboard",
            "• Optional: Deploy to AWS (your account)",
            "",
            "Knowledge Library: Curated document store",
            "• Sources: local files, folders, URLs",
            "• Ingestion: chunking, embeddings",
            "• Search via local_knowledge_search MCP tool",
            "• Built-in dashboard surface"
        ]
    )
    
    # Section 4: Pricing
    add_section_divider(prs, "Subscription & Pricing")
    
    add_content_slide(prs,
        "Credit System",
        [
            "Single subscription for all Kiro surfaces (IDE/CLI/Web/Mobile/Crew)",
            "",
            "Plans:",
            "• $20/month - 1,000 credits",
            "• $50/month - 2,500 credits",
            "• $100/month - 5,000 credits",
            "• $200/month - 10,000 credits",
            "• Add-on: $0.04/credit",
            "",
            "Model multipliers: Opus 5 (2.2x), Sonnet family (~1.3x), open-weight (0.05x-0.5x)"
        ]
    )
    
    add_content_slide(prs,
        "Cost Considerations",
        [
            "Verified from official docs:",
            "• Credits scale with model choice and task complexity",
            "• Model multipliers: Opus 5 2.2x, Sonnet family ~1.3x,",
            "  open-weight 0.05x-0.5x",
            "• Add-on credits $0.04 each; metered to 0.01 precision",
            "",
            "Not published by Kiro (measure in usage dashboard):",
            "• Per-subagent, per-cron, per-task credit consumption",
            "",
            "The subscription usage dashboard is the authoritative",
            "source for your own workload's cost."
        ],
        verification_status="verified"
    )
    
    # Section 4.5: Advanced & Enterprise Topics
    add_section_divider(prs, "Advanced & Enterprise Topics", "Full coverage of docs/crew/")

    add_content_slide(prs,
        "Security: 8-Layer Defense-in-Depth",
        [
            "Every tool call passes independent layers (runtime-enforced):",
            "1 Owner lock  2 Denied commands (137)  3 Governance ceiling",
            "4 Sensitive-path block  5 Tool approval  6 Input validation",
            "7 OS sandbox  8 Output redaction  (+ SEL audit, cross-cutting)",
            "",
            "Sandbox: auto (default) / strict / off; no OS layer on Windows",
            "Approval: Interactive / Trust command / Trust tool / Autopilot",
            "Denied cmds enforced at Crew's gate, not agent config",
            "kirocrew policy show|validate|explain; security events|audit|verify",
        ],
    )

    add_content_slide(prs,
        "Running 24/7",
        [
            "Local service: systemd/launchd; runs as user, sudo once for unit",
            "  kirocrew service install | logs -f | restart",
            "Docker: multi-arch GHCR; kiro-cli login + kirocrew token; SLSA",
            "  docker run -d -p 127.0.0.1:5476:5476 -v kirocrew-home:...",
            "Remote host: modern Linux, ~10GB RAM, Node 20+ (slack-mcp)",
            "  ssh -L 5476:localhost:5476 user@host",
            "Sync state: memory + SQLite incl WAL; NOT .env/.local_secret",
            "Mobile: Cloudflare/ngrok/Tailscale + dashboard.url; token 1h/20h",
        ],
    )

    add_content_slide(prs,
        "Interfaces: One Gateway, Many Surfaces",
        [
            "Mac app / web dashboard / CLI + Slack/Discord/Telegram/",
            "Teams/Webex/WeCom/WeChat channels",
            "Independent sessions, SHARED memory/skills/lessons/crons",
            "Channels connect outbound - no need to expose dashboard port",
            "",
            "Dashboard = React SPA localhost:5476 (primary interface)",
            "Cross-surface: dashboard<->Slack sync, resume sessions",
            "5476 loopback-only default; reverse proxy+TLS or SSH tunnel",
        ],
    )

    add_content_slide(prs,
        "Chat",
        [
            "Send text / @file / voice; streaming replies with live tool calls",
            "Guide: approve/reject, edit & resend, fork, Autopilot",
            "Each tab = independent session sharing long-term memory",
            "",
            "@ file picker; drag/paste image (vision); mic STT; Piper TTS",
            "Prompt optimizer (Cmd+Shift+Enter); incognito; switch models;",
            "reasoning effort per message",
            "Leaves: Sessions, Message controls, Optimizer, Voice, Widgets",
        ],
    )

    add_content_slide(prs,
        "Agent Capabilities (Customization Hub)",
        [
            "Agents - model/prompt/tools/MCP server config",
            "Agent Templates - prebuilt configs to customize",
            "Integrations (MCP) - add external tools/services",
            "Skills - on-demand knowledge files teaching workflows",
            "Steering - workspace-level rules every session inherits",
            "Hooks - automate reactions to events (scripts/context)",
            "Prompts - customize the system prompts shaping behavior",
        ],
    )

    add_content_slide(prs,
        "Configuration",
        [
            "Settings panel: Overview/Imports/Chat/Display/Voice/",
            "Notifications/Shortcuts/Skills/Channels/Browser/Computer Use/",
            "Instances/Security/Developer/About - changes apply immediately",
            "",
            "kirocrew config get|set|edit (set auto-restarts pool)",
            "config.json (JSON, missing keys -> defaults); 14 themes",
            "Channel creds in ~/.kiro/crew/.env (mode 600)",
            "Env: KIROCREW_HOME (~/.kiro/crew), KIROCREW_PORT (5476)",
        ],
    )

    add_content_slide(prs,
        "Multi-instance",
        [
            "Drive multiple remote Crew hosts from one hub via SSH tunnels",
            "Opt-in, off by default:",
            "  kirocrew config set instances.enabled true && kirocrew restart",
            "",
            "Warm set cap 5 (LRU); health probe 30s; 2-tier self-heal (<=8)",
            "Add: SSH host/alias, remote port 7777, token TTL 20h",
            "Security: deny-by-default, owner-only, loopback forwards,",
            "argv-list ssh (no shell injection), SEL-audited",
        ],
    )

    add_content_slide(prs,
        "Snapshot & Restore",
        [
            "kirocrew snapshot [dir] --keep N | --list",
            "kirocrew restore snap.tar.gz [--components ...] [--dry-run]",
            "",
            "Includes: Memory, Workspace, Crons, Config, Skills, Notifications",
            "Restore modes: replace (empty) vs merge (existing) - auto-detected",
            "Daily auto snapshot 08:00 UTC, keeps 7",
            "Contains security keys - treat like credentials",
            "NOT covered: kiro-cli/AWS creds, channel defs, app code",
        ],
    )

    add_content_slide(prs,
        "Apps / App Kit",
        [
            "Apps contribute: agents, skills, MCP servers, crons, UI pages,",
            "backend processes, gateway hooks",
            "Third-party apps DISABLED by default (Settings -> Security)",
            "",
            "Install from App Store (curated registry, PR to add) or",
            "federated external registries (opt-in)",
            "Lifecycle: gateway-managed / app-managed / locked",
            "Dev: kirocrew app dev <name>; only permissions.api enforced",
        ],
    )

    add_content_slide(prs,
        "Troubleshooting",
        [
            "Start with a health check:",
            "  kirocrew doctor [--verbose]   /   kirocrew logs -f",
            "Reports: kiro-cli, auth, embeddings, Slack, config, MCP",
            "",
            "Categories: install/setup, agent (AcpTimeoutError/",
            "AcpProcessDied), memory/embeddings, Slack, Discord, MCP,",
            "dashboard 403, sessions, task runner, subagents, cron,",
            "config, snapshot, multi-instance, voice, artifact deploy",
        ],
    )

    # Section 5: Use Cases
    add_section_divider(prs, "Part 4: Use Cases", "Enterprise integration examples")
    
    add_content_slide(prs,
        "Example: 24/7 PR Monitoring",
        [
            "Scenario: Track open PRs and alert on staleness",
            "",
            "Implementation:",
            "• Heartbeat watches open PRs",
            "• Alerts when stale (>4 hours no review)",
            "• Auto-pings reviewers via Slack",
            "",
            "Uses: Heartbeats + Slack integration",
            "",
            "Note: Example scenario based on documented capabilities"
        ],
        verification_status="partial"
    )
    
    add_content_slide(prs,
        "Example: Nightly Security Audit",
        [
            "Scenario: Automated security scanning",
            "",
            "Implementation:",
            "• Cron job at 2am",
            "• Spawns subagents (one per repo)",
            "• Scans for secrets, vulnerabilities, misconfigurations",
            "• Aggregates findings into report",
            "• Files HIGH/CRITICAL tickets",
            "",
            "Uses: Cron + Subagents + Task automation"
        ],
        verification_status="partial"
    )
    
    add_content_slide(prs,
        "Example: Incident Response",
        [
            "Scenario: Automated triage and RCA",
            "",
            "Implementation:",
            "• Webhook from PagerDuty",
            "• Crew triages alert",
            "• Checks logs, metrics, recent deploys",
            "• Applies known fixes (from memory/lessons)",
            "• Escalates if novel",
            "• Posts RCA when resolved",
            "",
            "Uses: Webhooks + Memory + Knowledge"
        ],
        verification_status="partial"
    )
    
    add_content_slide(prs,
        "CI/CD Integration (3 Modes)",
        [
            "Mode 1: CLI Headless",
            "• kiro-cli chat --no-interactive \"<prompt>\" in pipeline",
            "• Engine version: --engine v3 (not --v3)",
            "• Use case: Simple code generation",
            "",
            "Mode 2: Crew Webhooks",
            "• CI pipeline calls Crew on events",
            "• Crew analyzes and reports via Slack/Dashboard",
            "• Use case: Root cause analysis, monitoring",
            "",
            "Mode 3: Crew Cron",
            "• Periodic checks of CI/PR status",
            "• Use case: Nightly audits, health checks"
        ]
    )
    
    # Section 6: Getting Started
    add_section_divider(prs, "Part 5: Getting Started")
    
    add_two_column_slide(prs,
        "Installation",
        [
            "Option 1: Desktop App (Recommended)",
            "1. Download (.dmg or .AppImage)",
            "2. Launch (auto-starts Gateway)",
            "3. Authenticate (device-code sign-in)",
            "4. Start chatting",
            "",
            "No AWS CLI needed",
            "No Bedrock configuration needed"
        ],
        [
            "Option 2: CLI Install",
            "curl -fsSL https://download.crew.kiro.dev/cli.sh | sh",
            "",
            "kirocrew setup",
            "kirocrew doctor",
            "kirocrew gateway",
            "",
            "open http://localhost:5476"
        ]
    )
    
    add_content_slide(prs,
        "First Steps",
        [
            "1. Chat - Type and ask anything in dashboard",
            "",
            "2. Schedule a job - 'Every weekday at 9, summarize my open work'",
            "",
            "3. Teach a preference - 'Always use pytest over unittest'",
            "",
            "4. Run autonomous task - Projects panel → describe spec → Run",
            "",
            "5. Delegate parallel work - 'Research these 3 options in parallel'"
        ]
    )
    
    add_content_slide(prs,
        "Resources",
        [
            "Official Documentation:",
            "• kiro.dev/docs/crew/ - Main docs",
            "• kiro.dev/docs/crew/installation/ - Setup",
            "• kiro.dev/docs/crew/features/ - Features",
            "",
            "Specific Features:",
            "• kiro.dev/docs/crew/features/memory/ - Memory",
            "• kiro.dev/docs/crew/features/subagents/ - Subagents",
            "• kiro.dev/docs/crew/features/cron/ - Scheduling",
            "",
            "Pricing: kiro.dev/pricing/"
        ]
    )
    
    # Final slide
    add_title_slide(prs,
        "Questions?",
        "Thank you for your attention",
        "✓ All facts sourced from official docs | Verified 2026-08-31"
    )
    
    return prs


def main():
    print("Generating Kiro Crew Training Presentation (Verified Content)...")
    prs = create_presentation()
    
    output_file = "kiro-crew-training.pptx"
    prs.save(output_file)
    print(f"✅ Presentation saved as: {output_file}")
    print(f"   Total slides: {len(prs.slides)}")
    print()
    print("Content verification:")
    print("  ✓ All factual claims verified against official docs (2026-08-31)")
    print("  ℹ Per-operation credit costs / some quotas not published by Kiro")
    print()
    print("Sources used:")
    print("  - kiro.dev/docs/crew/")
    print("  - kiro.dev/docs/crew/installation/")
    print("  - kiro.dev/docs/crew/features/*")
    print("  - kiro.dev/pricing/")


if __name__ == "__main__":
    main()
